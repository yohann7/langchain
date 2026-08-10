"""Bounded PPTX parsing."""

from pathlib import Path


def load_pptx(path: Path, *, maximum_slides: int, maximum_shapes: int) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    if len(presentation.slides) > maximum_slides:
        raise ValueError("PPTX exceeds the configured slide limit")
    parts: list[str] = []
    shape_count = 0
    for number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {number}")
        for shape in slide.shapes:
            shape_count += 1
            if shape_count > maximum_shapes:
                raise ValueError("PPTX exceeds the configured shape limit")
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)
