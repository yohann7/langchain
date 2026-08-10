from pathlib import Path

import pytest

from knowledge_service.config import KnowledgeSettings
from knowledge_service.errors import UnsupportedDocumentError
from knowledge_service.parsing.chunking import split_document
from knowledge_service.parsing.redaction import redact_sensitive_content
from knowledge_service.parsing.registry import DocumentParser


def _parser(tmp_path: Path, **overrides) -> DocumentParser:
    return DocumentParser(
        KnowledgeSettings(
            run_dir=tmp_path / "runtime",
            embedding_device="cuda:0",
            **overrides,
        )
    )


def test_text_html_json_and_csv_are_normalized(tmp_path: Path) -> None:
    parser = _parser(tmp_path)
    text = tmp_path / "note.md"
    text.write_text("# 标题\r\n\r\n正文", encoding="utf-8")
    html = tmp_path / "page.html"
    html.write_text("<h1>Title</h1><script>bad()</script><p>Visible</p>", encoding="utf-8")
    data = tmp_path / "data.json"
    data.write_text('{"b": 2, "a": 1}', encoding="utf-8")
    table = tmp_path / "table.csv"
    table.write_text("name,value\nalpha,1\n", encoding="utf-8")

    assert parser.load(text).text == "# 标题\n\n正文"
    assert "bad()" not in parser.load(html).text
    assert parser.load(data).text.startswith("{\n  \"a\": 1")
    assert "name | value" in parser.load(table).text


def test_office_documents_are_parsed_by_format(tmp_path: Path) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    parser = _parser(tmp_path)
    docx_path = tmp_path / "report.docx"
    document = Document()
    document.add_heading("Summary", level=1)
    document.add_paragraph("Word body")
    document.save(docx_path)

    xlsx_path = tmp_path / "sheet.xlsx"
    workbook = Workbook()
    workbook.active.append(["name", "value"])
    workbook.active.append(["alpha", 1])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Deck title"
    slide.placeholders[1].text = "Slide body"
    presentation.save(pptx_path)

    assert "Word body" in parser.load(docx_path).text
    assert "alpha" in parser.load(xlsx_path).text
    assert "Slide body" in parser.load(pptx_path).text


def test_parser_rejects_unsupported_and_oversized_sources(tmp_path: Path) -> None:
    parser = _parser(tmp_path, maximum_source_bytes=1024)
    unsupported = tmp_path / "payload.exe"
    unsupported.write_bytes(b"x")
    too_large = tmp_path / "large.txt"
    too_large.write_bytes(b"x" * 1025)

    with pytest.raises(UnsupportedDocumentError):
        parser.load(unsupported)
    with pytest.raises(ValueError, match="size limit"):
        parser.load(too_large)


def test_pdf_parser_enforces_page_limit_and_extracts_page_markers(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    class Page:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class Reader:
        def __init__(self, _path) -> None:
            self.pages = [Page("first"), Page("second")]

    monkeypatch.setattr("pypdf.PdfReader", Reader)
    pdf = tmp_path / "manual.pdf"
    pdf.write_bytes(b"%PDF-test")

    with pytest.raises(ValueError, match="page limit"):
        _parser(tmp_path, maximum_pdf_pages=1).load(pdf)

    loaded = _parser(tmp_path, maximum_pdf_pages=2).load(pdf)
    assert "[Page 1]" in loaded.text
    assert "[Page 2]" in loaded.text


def test_redaction_and_chunking_preserve_useful_structure() -> None:
    redacted = redact_sensitive_content(
        "# Secrets\nDEEPSEEK_API_KEY=sk-abcdefghijklmnopqrstuvwxyz123456\nUseful text"
    )
    assert "sk-abcdefghijklmnopqrstuvwxyz123456" not in redacted.text
    assert redacted.redaction_count == 1
    assert redacted.was_redacted is True

    chunks = split_document(
        "# Heading\n" + "alpha beta gamma " * 40,
        source_type="markdown",
        chunk_size_chars=120,
        overlap_chars=20,
    )
    assert len(chunks) > 1
    assert chunks[0].heading_path == "Heading"
    assert all(chunk.content_hash for chunk in chunks)

    sections = split_document(
        "# First\nalpha\n# Second\nbeta",
        source_type="markdown",
        chunk_size_chars=200,
        overlap_chars=20,
    )
    assert [chunk.heading_path for chunk in sections] == ["First", "Second"]
